from base64 import b64encode
from magic import Magic
from fitz import open
from io import BytesIO 
from docx import Document
from docx.document import Document as _Document
from docx.oxml.text.paragraph import CT_P
from docx.oxml.table import CT_Tbl
from docx.table import _Cell, Table
from docx.text.paragraph import Paragraph
from pptx import Presentation
from hashlib import sha256
from pymongo import MongoClient
from datetime import datetime
from pytz import UTC
from os import path
from pptx.shapes.picture import Picture
from flask import Flask, request, jsonify
from spacy import load as spacy_load
from re import findall
from nltk import download
from nltk.corpus import stopwords
from nltk.stem import WordNetLemmatizer
app = Flask(__name__)
class MasterApp:
    def __init__(self):
        self.magic = Magic()

    def process_file(self, file_content):
        try:
            file_format = self.magic.from_buffer(file_content)
            content = {}
            print(file_format)

   
            if "pdf" in file_format.lower():
                content = self.process_pdf(file_content)
            elif "word" in file_format.lower():
                content = self.process_docx(file_content)
            elif "powerpoint" in file_format.lower():
                content = self.process_pptx(file_content)
            else:
                print("error: Unsupported file type.")
                return -1
            print("process file")
            return content
        except PermissionError:
            print("Permission denied")
            return 0
        except Exception as e:
            print(f"Error processing file: {e}")
            return -1
    
   
    def check_pdf_metadata(self,doc):
        metadata = doc.metadata
        for key, value in metadata.items():
            print(f"{key}: {value}")
        
        if "PowerPoint" in metadata.get('producer', '') or "PowerPoint" in metadata.get('creator', ''):
            return True
        else:
            return False
      
    def process_pdf(self, file):
        pdf_file = BytesIO(file)
        doc = open(stream=pdf_file, filetype="pdf")
        content = {}
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text().strip()
            image_list = page.get_images(full=True)
            images=[]
            for img in image_list:
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image_base64 = b64encode(image_bytes).decode('utf-8')

                images.append(image_base64)
                content[f"page_{page_num + 1}"] = {
                    "text": text,
                    "recognized_text": [],
                    "images": images,
                    "tables": []
                }  
                print(content,"\n")
        doc.close()
        return content

    def iter_block_items(self,parent):
        if isinstance(parent, _Document):
            parent_elm = parent.element.body
        elif isinstance(parent, _Cell):
            parent_elm = parent._tc
        else:
            raise ValueError("something's not right")

        for child in parent_elm.iterchildren():
            if isinstance(child, CT_P):
                yield Paragraph(child, parent)
            elif isinstance(child, CT_Tbl):
                yield Table(child, parent)

    def extract_table_content(self,table):
        table_data = []
        for row in table.rows:
            row_data = [cell.text.strip() for cell in row.cells]
            table_data.append(row_data)
        return table_data

    def process_docx(self, file):
        docx_file = BytesIO(file)
        doc = Document(docx_file)
        content = { 
            "text":"",
            "recognized_text": [],
            "images": [],
            "tables": []
        }
        # Extract text and tables
        for block in self.iter_block_items(doc):
            if isinstance(block, Paragraph):
                if(block.text.strip()!=""):
                    content["text"]+=(block.text.strip()+"\n")
            elif isinstance(block, Table):
                content["tables"].append(self.extract_table_content(block))

        # Extract images and perform OCR
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                image_part = rel.target_part
                image=image_part.blob
                image_base64 = b64encode(image).decode('utf-8')
                content["images"].append(image_base64)
                

        return {"page_1":content}

    def process_pptx(self, file):
        pptx_file = BytesIO(file)
        presentation = Presentation(pptx_file)
        content = {}
        tot_pg=0
        # Iterate over slides and extract text, images, and tables
        for slide in presentation.slides:
            tot_pg+=1
            text=""
            table_data=[]
            image_base64=[]
            ocr_text=[]
            # Extract text from the slide
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text_content = paragraph.text.strip()
                        if text_content:  # Only add non-empty text
                            text+=text_content
                    text+="\n"
                if shape.has_table:
                    table = shape.table
                    table_data.append(self.extract_table_content(table))  # Custom method to extract table content
                if isinstance(shape, Picture):
                    image = shape.image.blob
                    image_base64.append(b64encode(image).decode('utf-8'))

            content[f"page_{tot_pg}"] = {
                "text": text,
                "recognized_text": ocr_text,
                "images": image_base64,
                "tables": table_data
            }
        return content
class TextExtractor:
    def __init__(self):
        """
        Initialize text extraction utility with robust error handling
        """
        # Download necessary NLTK resources
        try:
            download('punkt', quiet=True)
            download('stopwords', quiet=True)
            download('wordnet', quiet=True)
        except:
            print("Warning: Unable to download all NLTK resources")
        
        self.nlp = spacy_load('en_core_web_sm')
        
        # Initialize lemmatizer and stopwords
        self.lemmatizer = WordNetLemmatizer()
        self.stop_words = set(stopwords.words('english'))

    def extract_key_information(self, text: str) -> dict:
        """
        Extract key information with categorized results, avoiding duplicates
        """
        # Process text with SpaCy
        doc = self.nlp(text)
        
        # Initialize results dictionary
        extracted_info = {
            'names': set(),
            'locations': set(),
            'organizations': set(),
            'dates': set(),
            'money': set(),
            'quantities': set(),
            'phone_numbers': set(),
            'emails': set(),
            'urls': set(),
            'other_phrases': set()
        }
        
        # Preserve multi-word named entities
        entity_labels = {
            'ORG': 'organizations',
            'GPE': 'locations',
            'LOC': 'locations',
            'PERSON': 'names',
            'MONEY': 'money',
            'QUANTITY': 'quantities'
        }
        
        # Track all extracted values to prevent duplicates
        all_extracted_values = set()
        
        # Extract named entities
        for ent in doc.ents:
            if ent.label_ in entity_labels:
                category = entity_labels[ent.label_]
                value = ent.text.strip()
                if value not in all_extracted_values:
                    extracted_info[category].add(value)
                    all_extracted_values.add(value)
        
        # Advanced key phrase extraction
        def extract_meaningful_phrases(doc):
            phrases = []
            current_phrase = []
            for token in doc:
                if token.pos_ in ['NOUN', 'PROPN', 'ADJ']:
                    current_phrase.append(token.text)
                else:
                    if len(current_phrase) > 1:
                        phrases.append(' '.join(current_phrase))
                    current_phrase = []
            
            if len(current_phrase) > 1:
                phrases.append(' '.join(current_phrase))
            
            return phrases
        
        # Extract meaningful multi-word phrases
        meaningful_phrases = extract_meaningful_phrases(doc)
        for phrase in meaningful_phrases:
            if phrase not in all_extracted_values:
                extracted_info['other_phrases'].add(phrase)
                all_extracted_values.add(phrase)
        
        # Date extraction
        date_patterns = [
            r'\b\d{1,2}[/.-]\d{1,2}[/.-]\d{4}\b',  # DD/MM/YYYY or MM/DD/YYYY
            r'\b\d{4}[/.-]\d{1,2}[/.-]\d{1,2}\b',  # YYYY/MM/DD
            r'\b(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)\s+\d{1,2},\s+\d{4}\b'  # Month DD, YYYY
        ]
        
        # Extract dates
        for pattern in date_patterns:
            dates = set(findall(pattern, text))
            for date in dates:
                if date not in all_extracted_values:
                    extracted_info['dates'].add(date)
                    all_extracted_values.add(date)
        
        # Email extraction
        email_pattern = r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b'
        for email in findall(email_pattern, text):
            if email not in all_extracted_values:
                extracted_info['emails'].add(email)
                all_extracted_values.add(email)
        
        # Phone number extraction
        phone_patterns = [
            r'\b(?:\+1\s?)?(?:\(\d{3}\)|\d{3})[-.\s]?\d{3}[-.\s]?\d{4}\b',  # US/International
            r'\b\d{10}\b',  # 10 digit numbers without formatting
            r'\b\(\d{3}\)\s?\d{3}[-.]?\d{4}\b'  # (XXX) XXX-XXXX format
        ]
        for pattern in phone_patterns:
            found_numbers = findall(pattern, text)
            for num in found_numbers:
                # Handle both tuple and string results
                if isinstance(num, tuple):
                    cleaned_num = ''.join(filter(bool, num))
                    if cleaned_num and cleaned_num not in all_extracted_values:
                        extracted_info['phone_numbers'].add(cleaned_num)
                        all_extracted_values.add(cleaned_num)
                elif isinstance(num, str):
                    if num not in all_extracted_values:
                        extracted_info['phone_numbers'].add(num)
                        all_extracted_values.add(num)
        
        # URL extraction
        url_pattern = r'https?://\S+|www\.\S+'
        for url in findall(url_pattern, text):
            if url not in all_extracted_values:
                extracted_info['urls'].add(url)
                all_extracted_values.add(url)
        
        # Remove empty sets
        extracted_info = {k: v for k, v in extracted_info.items() if v}
        
        return extracted_info
client = MongoClient("mongodb+srv://jeevajoslin:p7MK68VRoY7LvGXh@tratos.lt7g3.mongodb.net/")
db = client["TRANS2"]  
trans_collection = db["TRANS"] 
userfiles_collection = db["USERFILES"]
def generate_sha256(file):
    return sha256(file).hexdigest()
@app.route('/process', methods=['POST'])
def upload_file():
    app.master_app = MasterApp()
    if 'file' not in request.files:
        return jsonify({"error": "No file provided"}), 400
    file = request.files['file']
    print(request.files)
    user = request.form.get('user')
    print(user)
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400
    file_content = file.read()

    userfiles_document = userfiles_collection.find_one({"username":user})
    print(userfiles_document)
    if not userfiles_document:
        return jsonify({"error": "User not found. Cannot upload files."}), 404
    sha_key = generate_sha256(file_content)
    print(sha_key)
    existing_document = trans_collection.find_one({"sha_key": sha_key})
    if existing_document:
            # Consistent return type
            userfiles_collection.update_one(
                {"username": user},
                {"$addToSet": {"sha_keys": sha_key}} 
            )
            return jsonify({
                "message": "File already exists in the database",
                "sha_key": sha_key,
                "existing_content": existing_document.get("json_file_content", {})
            }), 200
    content = app.master_app.process_file(file_content)
    if content!=-1 and content!=0:
        text=""
        for _, page_content in content.items():
            # Extract recognized text if available
            if 'recognised_text' in page_content:
                text+=" ".join(page_content["recognised_text"])
            
            # Extract regular text if available
            if 'text' in page_content:
                text+=(" "+page_content['text'])
        extract=TextExtractor().extract_key_information(text)
        meta={}
        for category, values in extract.items():
            meta[category.replace('_', ' ')]=list(values)
        print("meta",meta)
        content["metadata"]=meta
    print("last",jsonify(content))
    print(content)
    if content == -1:
        return jsonify({"error": "Unsupported file type"}), 400
    elif content ==0:
        return jsonify({"error": "Permission denied"}), 400
    
    document = {
        "file_name": path.basename(file.filename),
        "file_type": path.splitext(file.filename)[1],
        "sha_key": sha_key,
        "processed_at": datetime.now(UTC)
    }
    
    result = trans_collection.insert_one(document)
    trans_collection.update_one(
        {"_id": result.inserted_id},
        {"$set": {"json_file_content": content}}
    )

    userfiles_collection.update_one(
        {"username": user},
        {"$addToSet": {"sha_keys": sha_key}} 
    )
    return jsonify(content), 200
app.run(debug=True)
